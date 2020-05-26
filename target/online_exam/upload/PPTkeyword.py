# -*- coding: utf-8 -*
from pptx.api import Presentation
from collections import Counter
import jieba
import sys
import os

filepath = '/home/xzq/桌面/1.pptx'
# 停用词的路径
baidu_stopwords_filepath = '/home/xzq/桌面/stopWord/stopwords/baidu_stopwords.txt'
cn_stopwords_filepath = '/home/xzq/桌面/stopWord/stopwords/cn_stopwords.txt'
hit_stopwords_filepath = '/home/xzq/桌面/stopWord/stopwords/hit_stopwords.txt'
scu_stopwords_filepath = '/home/xzq/桌面/stopWord/stopwords/scu_stopwords.txt'
stopWords_filePathList = [
    baidu_stopwords_filepath, cn_stopwords_filepath, hit_stopwords_filepath, scu_stopwords_filepath
]

def stopwordslist(stopwordsDir):  # 利用得到的停用词文件获取停用词列表
    filePathList = os.listdir(stopwordsDir)
    # print(filePathList)
    stopWordsTotal = set()
    for each_Path in filePathList:
        filepath = stopwordsDir + '/' + each_Path
        each_Path_Words = [line.strip() for line in open(filepath, 'r').readlines()]
        for eachWord in each_Path_Words:
            stopWordsTotal.add(eachWord)
    return list(stopWordsTotal)

# stopwordslist('/home/xzq/桌面/stopWord')
# 对句子进行分词
def seg_sentence(sentence, stopwordsDir):
    sentence_seged = jieba.cut(sentence.strip())
    stopwords = stopwordslist(stopwordsDir)  # 加载停用词的路径
    outstr = ''
    for word in sentence_seged:
        if word not in stopwords:
            if word != '\t':
                outstr += word
                outstr += " "
    return outstr



# result = stopwordslist(stopWords_filePathList)
# print(result)
# print("finalNum:", len(list(stopWordsTotal)))

def get(stopwordsDir, filepath, pageNums=None):

    # 按照pageNums中PPT页码识别关键词，若为None,则遍历提取整个PPT文件

    result = ""
    ppt = Presentation(pptx=filepath)
    # print(len(ppt.slides))
    slides = ppt.slides
    if pageNums == None:
        pageNums = range(1, len(slides)+1)
    else:
        pageNums = [int(each) for each in pageNums]
    # else:
    #     slides = [ppt.slides[pageNum] for pageNum in pageNums]
    # slides = ppt.slides
    # page = 1
    # print(pageNums)
    for page in pageNums:
        slide = slides[page-1]
        text_runs = []
        s = ""          #　该页ppt上所有的文字
        text_types = []     #　字体类型
        text_sizes = []     #  字体大小
        text_rgb = []       #  字体颜色
        text_bold = []      #  是否黑体
        text_italic = []    #  是否斜体
        text_underline = [] #  是否下划线
        rates = {}
        shapes = slide.shapes

        for shape in shapes:
            if type(shape).__name__ == 'Picture':   # BaseShape的类型为‘Picture’
                pass

            if type(shape).__name__ == 'GraphicFrame':   # BaseShape 的类型为'GraphicFrame'
                pass

            else:   # BaseShape的类型为‘Shape’
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                paragraphs = tf.paragraphs
                for paragraph in paragraphs:

                    runs = paragraph.runs
                    for run in runs:
                        text = run.text
                        # if text[0] == ' ':  # 空字符过滤
                        #     continue
                        if str(run.font.color.type) == 'None':  # 默认颜色为黑色
                            text_rgb.append('000000')
                        else:
                            text_rgb.append(run.font.color.rgb)

                        if str(run.font.size) == 'None':   # 默认字体大小为18
                            text_sizes.append(18.0)
                        else:
                            text_sizes.append(run.font.size.pt)
                        # print(str(run.font.size))

                        if str(run.font.bold) == 'None':  # 粗体
                            text_bold.append(False)
                        else:
                            text_bold.append(run.font.bold)

                        # print(run.font.bold)
                        if str(run.font.italic) == 'None':  # 斜体
                            text_italic.append(False)
                        else:
                            text_italic.append(run.font.italic)


                        if str(run.font.underline) == 'None':  # 下划线
                            text_underline.append(False)
                        else:
                            text_underline.append(run.font.underline)
                        # text_italic.append(run.font.italic)

                        if str(run.font.name) == 'None':  # 默认为宋体
                            text_types.append('宋体')
                        else:
                            text_types.append(run.font.name)  # 字体名

                        # print(run.font.name)
                        s += text
                        text_runs.append(run.text)

        text_types_dict = Counter(text_types)
        text_sizes_dict = Counter(text_sizes)
        text_rgb_dict = Counter(text_rgb)
        text_bold_dict = Counter(text_bold)
        text_italic_dict = Counter(text_italic)
        text_underline_dict = Counter(text_underline)
        rgb_count = len(text_rgb)
        size_count = len(text_sizes)
        type_count = len(text_types)
        bold_count = len(text_bold)
        italic_count = len(text_italic)
        underline_count = len(text_underline)

        goal = []

        for shape in shapes:
            if type(shape).__name__ == 'Picture':  # BaseShape的类型为‘Picture’
                pass

            elif type(shape).__name__ == 'GraphicFrame':  # BaseShape 的类型为'GraphicFrame'
                pass

            else:  # BaseShape的类型为‘Shape’
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                paragraphs = tf.paragraphs
                for paragraph in paragraphs:
                    # print(type(paragraph))
                    runs = paragraph.runs
                    for run in runs:
                        text = run.text
                        if text[0] == ' ':  # 空字符过滤
                            continue

                        if str(run.font.color.type) == 'None':  # 默认颜色为黑色
                            rgb_rate = 1 - text_rgb_dict['000000'] / rgb_count  # 颜色比重
                        else:
                            rgb_rate = 1 - text_rgb_dict[run.font.color.rgb] / rgb_count # 颜色比重

                        if str(run.font.size) == 'None':   # 默认字体大小为18
                            size_rate = 1 - text_sizes_dict['18.0'] / size_count  # 字体比重
                        else:
                            size_rate = 1 - text_sizes_dict[run.font.size.pt] / size_count

                        if str(run.font.name) == 'None':
                            type_rate = 1 - text_types_dict['宋体'] / type_count
                        else:
                            # text_types.append(run.font.name) # 字体名
                            type_rate = 1 - text_types_dict[run.font.name] / type_count

                        if str(run.font.bold) == 'None':  # 粗体
                            bold_rate = 1 - text_bold_dict[False] / bold_count
                            # text_bold.append(False)
                        else:
                            # text_bold.append(run.font.bold)
                            bold_rate = 1 - text_bold_dict[run.font.bold] / bold_count
                        # print(run.font.bold)

                        if str(run.font.italic) == 'None': # 斜体
                            # text_italic.append(False)
                            italic_rate = 1 - text_italic_dict[False] / italic_count
                        else:
                            italic_rate = 1 - text_italic_dict[run.font.italic] / italic_count
                            # text_italic.append(run.font.italic)


                        if str(run.font.underline) == 'None':  #下划线
                            # text_underline.append(False)
                            underline_rate = 1 - text_underline_dict[False] / underline_count
                        else:
                            # text_underline.append(run.font.underline)
                            underline_rate = 1 - text_underline_dict[run.font.underline] / underline_count
                        # text_italic.append(run.font.italic)
                        if text in rates.keys():
                            curSum = sum([rgb_rate, size_rate, type_rate, bold_rate, italic_rate, underline_rate])
                            if sum(rates[text]) < curSum:
                                rates[text] = [rgb_rate, size_rate, type_rate, bold_rate, italic_rate, underline_rate]
                        else:
                            rates[text] = [rgb_rate, size_rate, type_rate, bold_rate, italic_rate, underline_rate]
                        goal.append(20 * (rgb_rate + size_rate + type_rate + bold_rate + italic_rate + underline_rate))

                        s += text

        # print(rates)
        text_runs.append(run.text)
        # print(text_runs)
        # print(s)
        key_word = text_runs[goal.index(max(goal))]
        # print(key_word)
        # print(key_word)
        seg = seg_sentence(key_word, stopwordsDir)
        # print("Default Mode: " + "/ ".join(seglist))
        # print(seglist)
        # result=""
        result += "第" + str(page) + "页关键字为:" + seg + "\n"
        # print("第%d页 关键字为:%s"%(page,seg))
        page += 1
    # print(result)

    return result


# get(filepath, ['1', '2', '3'])
# print(stopwordslist(stopWords_filePathList))
#
if __name__ == '__main__':
    params = []
    for i in range(1, len(sys.argv)):
        params.append((sys.argv[i]))
    # print(params)
    if len(params) == 1:
        print(get(params[0]))
    # print(len(stopwordslist(params)))
    # print(stopwordslist(params))
    # print(sys.argv[0])
    else:
        pageNums = params[2].split(',')
        # print(pageNums)
        print(get(params[0], params[1], pageNums))

# get(filepath)
# print(key_word)
# print(goal)
# print(rates)
# print(text_runs)
# print(s)



